"""Multi-format document-loader registry tests (P1: multi-format ingestion).

Covers the pluggable loader registry (`loaders.py`) and its integration with the
existing chunk → embed → index pipeline in `custom_rag.upload`:

- Per-format plain-text extraction with tiny in-memory fixtures
  (DOCX/PPTX via python-docx/python-pptx, HTML via BeautifulSoup, MD, TXT).
- Unknown / missing extensions raise `UnsupportedFormatError` (the 4xx/skip
  backstop) — never a crash deeper in the pipeline.
- Malformed bytes for a *recognized* format raise `ExtractionError`.
- The registry helpers (`get_extension`, `is_supported_name`,
  `is_registry_format`, `SUPPORTED_EXTENSIONS`).
- An end-to-end `upload()` of a non-PDF document (read mode) reaches the index
  via the same chunk/embed/upsert path PDFs use — confirming the registry path
  produces real chunks, mints the deterministic id scheme, and never touches the
  PDF-only fitz / Document-Intelligence code.

All tests are hermetic — no network, no real Azure SDK calls, no PDFs.
"""

import io
import os
import sys

import pytest

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

# Azure env only needs to be present so module import / client construction does
# not blow up; every Azure call is mocked in the integration test below.
os.environ.setdefault("AZURE_OPENAI_ENDPOINT", "https://test.openai.azure.com/")
os.environ.setdefault("AZURE_OPENAI_KEY", "test-key")
os.environ.setdefault("AZURE_OPENAI_VERSION", "2024-02-01")
os.environ.setdefault("AZURE_OPENAI_EMBEDDING_MODEL", "text-embedding-3-small")
os.environ.setdefault("AZURE_SEARCH_ENDPOINT", "https://test.search.windows.net")
os.environ.setdefault("AZURE_SEARCH_KEY", "test-key")
os.environ.setdefault("INDEX_NAME", "test-index")
os.environ.setdefault("DOC_INTELLIGENCE_ENDPOINT", "https://test.cognitiveservices.azure.com/")
os.environ.setdefault("DOC_INTELLIGENCE_KEY", "test-key")

from loaders import (  # noqa: E402
    PDF_EXTENSION,
    SUPPORTED_EXTENSIONS,
    ExtractionError,
    UnsupportedFormatError,
    extract_text,
    get_extension,
    is_registry_format,
    is_supported_name,
)

pytestmark = pytest.mark.ingestion


# ---------------------------------------------------------------------------
# Tiny in-memory fixtures (no fixture files on disk — generated per test)
# ---------------------------------------------------------------------------


def _docx_bytes(paragraphs: list[str]) -> bytes:
    """Build a minimal valid DOCX (OOXML) in memory via python-docx."""
    from docx import Document

    doc = Document()
    for p in paragraphs:
        doc.add_paragraph(p)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _pptx_bytes(slides: list[list[str]]) -> bytes:
    """Build a minimal valid PPTX deck in memory via python-pptx.

    `slides` is a list of slides; each slide is a list of text lines added to a
    single textbox on a blank layout.
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[6]  # blank layout
    for lines in slides:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
        tf = box.text_frame
        tf.text = lines[0] if lines else ""
        for extra in lines[1:]:
            tf.add_paragraph().text = extra
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Per-format extraction
# ---------------------------------------------------------------------------


def test_extract_txt():
    doc = extract_text("notes.txt", b"hello world\nsecond line")
    assert doc.parser == "text"
    assert doc.page_count == 1
    assert "hello world" in doc.text
    assert "second line" in doc.text


def test_extract_txt_lenient_decode():
    # Undecodable bytes are replaced, not fatal.
    doc = extract_text("weird.txt", b"caf\xe9 bytes")
    assert "caf" in doc.text  # the valid prefix survives the lenient decode


def test_extract_markdown_keeps_headers():
    md = b"# Title\n\nBody paragraph with content.\n\n## Section\n\nMore text."
    doc = extract_text("guide.md", md)
    assert doc.parser == "markdown"
    # Markdown is returned verbatim so layout-mode header splitting still works.
    assert "# Title" in doc.text
    assert "## Section" in doc.text


def test_extract_html_strips_scripts_and_styles():
    html = (
        b"<html><head><style>.x{color:red}</style></head>"
        b"<body><h1>Heading</h1>"
        b"<script>var secret='do-not-index';</script>"
        b"<p>Visible body text.</p></body></html>"
    )
    doc = extract_text("page.html", html)
    assert doc.parser == "html"
    assert "Heading" in doc.text
    assert "Visible body text." in doc.text
    # script / style content must not reach the index.
    assert "do-not-index" not in doc.text
    assert "color:red" not in doc.text


def test_extract_htm_alias():
    doc = extract_text("legacy.htm", b"<html><body><p>Legacy page.</p></body></html>")
    assert doc.parser == "html"
    assert "Legacy page." in doc.text


def test_extract_docx():
    content = _docx_bytes(["First paragraph.", "", "Second paragraph."])
    doc = extract_text("report.docx", content)
    assert doc.parser == "docx"
    assert doc.page_count == 1
    assert "First paragraph." in doc.text
    assert "Second paragraph." in doc.text


def test_extract_pptx_counts_slides():
    content = _pptx_bytes([["Slide one title", "bullet a"], ["Slide two title"]])
    doc = extract_text("deck.pptx", content)
    assert doc.parser == "pptx"
    assert doc.page_count == 2  # slide count
    assert "Slide one title" in doc.text
    assert "Slide two title" in doc.text


# ---------------------------------------------------------------------------
# Unknown extension + malformed bytes
# ---------------------------------------------------------------------------


def test_unknown_extension_raises_unsupported():
    with pytest.raises(UnsupportedFormatError):
        extract_text("archive.zip", b"PK\x03\x04 not really handled")


def test_missing_extension_raises_unsupported():
    with pytest.raises(UnsupportedFormatError):
        extract_text("README", b"no extension here")


def test_unsupported_error_message_omits_filename():
    # Safe messages: the error names the extension, never the (possibly
    # path-bearing) filename.
    with pytest.raises(UnsupportedFormatError) as exc:
        extract_text("/secret/path/file.xyz", b"data")
    assert "/secret/path" not in str(exc.value)
    assert ".xyz" in str(exc.value)


def test_pdf_not_handled_by_registry():
    # PDF text extraction stays in upload(); the registry rejects it as a
    # defense-in-depth backstop (callers route .pdf to the fitz/DI path).
    with pytest.raises(UnsupportedFormatError):
        extract_text("doc.pdf", b"%PDF-1.7 fake")


def test_malformed_docx_raises_extraction_error():
    with pytest.raises(ExtractionError):
        extract_text("broken.docx", b"this is not a real OOXML zip")


def test_malformed_pptx_raises_extraction_error():
    with pytest.raises(ExtractionError):
        extract_text("broken.pptx", b"not a zip at all")


# ---------------------------------------------------------------------------
# Registry helpers
# ---------------------------------------------------------------------------


def test_get_extension_uses_last_dot_and_lowercases():
    assert get_extension("report.final.DOCX") == ".docx"
    assert get_extension("path/to/notes.TXT") == ".txt"
    assert get_extension("noext") == ""
    assert get_extension(".dotfile") == ""  # leading-dot dotfile, no real ext


def test_supported_extensions_set():
    assert PDF_EXTENSION in SUPPORTED_EXTENSIONS
    for ext in (".docx", ".pptx", ".html", ".htm", ".md", ".txt"):
        assert ext in SUPPORTED_EXTENSIONS


def test_is_supported_name():
    assert is_supported_name("a.pdf")
    assert is_supported_name("a.DOCX")
    assert is_supported_name("a.md")
    assert not is_supported_name("a.png")
    assert not is_supported_name("noext")


def test_is_registry_format_excludes_pdf():
    # PDF is "supported" but NOT a registry format (it keeps the fitz/DI path).
    assert is_supported_name("a.pdf")
    assert not is_registry_format("a.pdf")
    assert is_registry_format("a.docx")
    assert is_registry_format("a.txt")


# ---------------------------------------------------------------------------
# Integration: a non-PDF flows through upload() (read mode) to the index
# ---------------------------------------------------------------------------


class _FakeUploadFile:
    def __init__(self, content: bytes, filename: str):
        self._content = content
        self.filename = filename

    async def read(self) -> bytes:
        return self._content


class _FakeSearchClient:
    def __init__(self):
        self.uploaded = None

    def search(self, **kwargs):
        return []  # no stale chunks

    def delete_documents(self, documents):
        return None

    def merge_or_upload_documents(self, documents):
        self.uploaded = documents

        class _Result:
            succeeded = True
            key = "k"

        return [_Result() for _ in documents]


def _patched_rag(monkeypatch):
    """Build a rag() with search + embeddings + token counting mocked.

    Deliberately does NOT mock fitz or AzureAIDocumentIntelligenceLoader: the
    registry path must never touch them, so leaving them unmocked is part of the
    assertion — if the code wrongly routed a non-PDF through the PDF branch, the
    fake bytes would blow up there.
    """
    from unittest.mock import AsyncMock

    from custom_rag import rag

    instance = rag()
    fake_search = _FakeSearchClient()
    monkeypatch.setattr(instance, "create_search_index", AsyncMock(return_value=fake_search))
    monkeypatch.setattr(instance, "get_embedding", AsyncMock(return_value=[0.1] * 3))
    monkeypatch.setattr(instance, "chunk_token", AsyncMock(return_value=42))
    return instance, fake_search


@pytest.mark.asyncio
async def test_upload_txt_read_mode_indexes_chunks(monkeypatch):
    instance, fake_search = _patched_rag(monkeypatch)
    body = ("The quick brown fox jumps over the lazy dog. " * 40).encode("utf-8")
    file = _FakeUploadFile(body, "notes.txt")

    result = await instance.upload(file, "tenant1", "read", "/srv/notes.txt")

    # Real chunks were produced and upserted via the shared pipeline.
    assert isinstance(result, dict)
    assert result["status"] == "successful"
    assert result["total_chunks"] >= 1
    assert fake_search.uploaded is not None and len(fake_search.uploaded) >= 1

    # Deterministic id scheme (P0-4) and provenance are unchanged for non-PDFs.
    first = fake_search.uploaded[0]
    assert first["id"].startswith("tenant1_")
    assert "_read_" in first["id"]
    assert first["bot_tag"] == "tenant1"
    assert first["source_path"] == "/srv/notes.txt"
    assert first["filename"] == "notes.txt"


@pytest.mark.asyncio
async def test_upload_docx_read_mode_indexes_chunks(monkeypatch):
    instance, fake_search = _patched_rag(monkeypatch)
    content = _docx_bytes(["A paragraph of real content. " * 20, "Another paragraph here. " * 20])
    file = _FakeUploadFile(content, "report.docx")

    result = await instance.upload(file, "tenant1", "read", "/srv/report.docx")

    assert isinstance(result, dict)
    assert result["status"] == "successful"
    assert result["total_chunks"] >= 1
    assert fake_search.uploaded and fake_search.uploaded[0]["id"].startswith("tenant1_")


@pytest.mark.asyncio
async def test_upload_markdown_layout_mode_splits_on_headers(monkeypatch):
    instance, fake_search = _patched_rag(monkeypatch)
    md = (
        b"# Intro\n\nIntro body text with enough words to chunk.\n\n"
        b"## Details\n\nDetailed section body text here.\n"
    )
    file = _FakeUploadFile(md, "guide.md")

    # .md is the one new format that benefits from layout (header) mode.
    result = await instance.upload(file, "tenant1", "layout", "/srv/guide.md")

    assert isinstance(result, dict)
    assert result["status"] == "successful"
    assert result["total_chunks"] >= 1
    # Header splitting captured the markdown section header.
    section_headers = {d.get("section_header", "") for d in fake_search.uploaded}
    assert "Intro" in section_headers


@pytest.mark.asyncio
async def test_upload_unknown_extension_raises_from_registry(monkeypatch):
    # Defense-in-depth: even if the route guard were bypassed, upload() routes an
    # unknown extension to the registry, which raises UnsupportedFormatError
    # (a 4xx-friendly error) rather than crashing in the PDF branch.
    instance, _ = _patched_rag(monkeypatch)
    file = _FakeUploadFile(b"some bytes", "mystery.xyz")
    with pytest.raises(UnsupportedFormatError):
        await instance.upload(file, "tenant1", "read", "/srv/mystery.xyz")
