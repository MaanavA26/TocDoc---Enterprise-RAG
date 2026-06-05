"""Pluggable document-loader registry (multi-format ingestion).

The ingestion write path — ``custom_rag.rag.upload`` — was historically PDF-only:
it analyzed the bytes with PyMuPDF (``fitz``) and extracted text via Azure
Document Intelligence. This module adds a small, dependency-light registry that
dispatches by **file extension** to extract PLAIN TEXT from a handful of common
office/web/text formats, so the EXISTING chunk → embed → index pipeline can
ingest them unchanged.

Design constraints (deliberate, to keep the blast radius tiny):

- **Extraction only.** A loader returns plain text plus a coarse page/slide
  count and a parser name for observability. It NEVER chunks, embeds, mints
  chunk IDs, or touches the index — those stay in ``upload()`` so the P0-4
  deterministic-ID and P0-5 chunking guarantees remain enforced in one place.
- **PDF is NOT handled here.** ``upload()`` keeps its exact PyMuPDF + Document
  Intelligence path for ``.pdf``. This registry covers only the new formats.
- **Unknown extensions raise ``UnsupportedFormatError``** (a 4xx-friendly,
  P0-6-safe error) rather than crashing deeper in the pipeline. Callers
  (the ``/upload`` route, connector ``enumerate``) pre-filter on
  ``SUPPORTED_EXTENSIONS`` so an unknown type is a clean skip / 415, never a 500.
- **Safe messages only.** Error text carries the extension/class, never file
  bytes or document content.
"""

from __future__ import annotations

import io
import logging
import zipfile
from collections.abc import Callable

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Zip-bomb / decompression-bomb guard (OOXML .docx/.pptx)
# ---------------------------------------------------------------------------
#
# .docx and .pptx are zip+XML containers. python-docx / python-pptx decompress
# the document parts fully into memory and hand multi-GB of XML to lxml. A
# crafted OOXML file well under the 100 MB *compressed* upload ceiling can
# decompress ~1000x to several GB of XML and OOM the worker. The only size
# ceiling elsewhere in the ingestion path is enforced on the COMPRESSED bytes
# (app.py / connectors), so we bound the UNCOMPRESSED size and the
# decompression ratio here — the one place both the /upload route and the
# (non-admin-authored) connector ingest path funnel through.
#
# Module-level constants so tests can monkeypatch a small cap and use a tiny,
# highly-compressible fixture instead of fabricating multi-GB inputs.
_MAX_OOXML_UNCOMPRESSED_BYTES = 500 * 1024 * 1024  # 500 MB total inflated size
_MAX_OOXML_COMPRESSION_RATIO = 100  # uncompressed / compressed


def _guard_ooxml_zip_bomb(content: bytes, fmt: str) -> None:
    """Reject an OOXML payload that would decompress beyond safe bounds.

    Opens ``content`` as a zip and sums every member's declared uncompressed
    size BEFORE python-docx/pptx inflates it. Raises ``ExtractionError`` (a
    4xx-friendly loader error) when the total inflated size or the
    uncompressed/compressed ratio exceeds the configured cap, or when the bytes
    are not a valid zip at all. The message names the format and the bound only,
    never document content.

    Args:
        content: Raw OOXML (zip container) bytes.
        fmt: Short format label for the error message (e.g. "DOCX").
    """
    try:
        with zipfile.ZipFile(io.BytesIO(content)) as zf:
            total_uncompressed = sum(info.file_size for info in zf.infolist())
    except zipfile.BadZipFile as exc:
        # Not a valid OOXML/zip container — surface as a recognized-but-malformed
        # error so the route maps it to a 4xx (not a 500), same as a corrupt doc.
        raise ExtractionError(f"Failed to parse {fmt} (BadZipFile)") from exc

    if total_uncompressed > _MAX_OOXML_UNCOMPRESSED_BYTES:
        raise ExtractionError(
            f"{fmt} rejected: uncompressed size exceeds the "
            f"{_MAX_OOXML_UNCOMPRESSED_BYTES // (1024 * 1024)} MB limit."
        )

    compressed = len(content)
    if compressed > 0 and total_uncompressed / compressed > _MAX_OOXML_COMPRESSION_RATIO:
        raise ExtractionError(f"{fmt} rejected: decompression ratio exceeds {_MAX_OOXML_COMPRESSION_RATIO}x.")


# ---------------------------------------------------------------------------
# Errors
# ---------------------------------------------------------------------------


class LoaderError(Exception):
    """Base class for document-loader failures (P0-6 friendly, safe messages)."""


class UnsupportedFormatError(LoaderError):
    """The file extension has no registered loader.

    Raised by ``extract_text`` for an extension outside ``SUPPORTED_EXTENSIONS``
    (and outside the PDF path, which the registry does not own). Callers map
    this to a 4xx / skip — it must never surface as a 500.
    """


class ExtractionError(LoaderError):
    """A registered loader failed to parse otherwise-recognized bytes.

    e.g. a ``.docx`` whose bytes are not a valid OOXML zip. The message names
    the format and the underlying error CLASS only — never raw content.
    """


# ---------------------------------------------------------------------------
# Extraction result
# ---------------------------------------------------------------------------


class ExtractedDocument:
    """Plain-text payload handed back to ``upload()``.

    Mirrors only what the downstream chunking branches read: ``text`` (fed in as
    ``docs[0].page_content``), ``page_count`` (reported as ``total_pages`` /
    ``page_count`` in observability), and ``parser`` (the ``document_parsed``
    event's parser name). No content is logged.
    """

    __slots__ = ("text", "page_count", "parser")

    def __init__(self, *, text: str, page_count: int, parser: str) -> None:
        self.text = text
        self.page_count = page_count
        self.parser = parser


# ---------------------------------------------------------------------------
# Individual format loaders (extension → bytes → ExtractedDocument)
# ---------------------------------------------------------------------------


def _load_txt(content: bytes) -> ExtractedDocument:
    """Decode a plain-text file as UTF-8 (lenient: undecodable bytes replaced)."""
    text = content.decode("utf-8", errors="replace")
    return ExtractedDocument(text=text, page_count=1, parser="text")


def _load_markdown(content: bytes) -> ExtractedDocument:
    """Decode Markdown as UTF-8 text.

    The raw Markdown (headers and all) is returned verbatim so that ``layout``
    mode's ``MarkdownHeaderTextSplitter`` can split on the real ``#`` headers —
    .md is the one new format that genuinely benefits from layout mode.
    """
    text = content.decode("utf-8", errors="replace")
    return ExtractedDocument(text=text, page_count=1, parser="markdown")


def _load_html(content: bytes) -> ExtractedDocument:
    """Extract visible text from HTML using BeautifulSoup's stdlib parser.

    Uses the built-in ``html.parser`` so no hard ``lxml`` *parser* dependency is
    added (lxml still arrives transitively via python-docx/pptx, but the HTML
    path does not require it). ``<script>``/``<style>``/``<head>`` are dropped so
    only human-visible body text is indexed.
    """
    from bs4 import BeautifulSoup

    try:
        soup = BeautifulSoup(content, "html.parser")
    except Exception as exc:  # noqa: BLE001 - normalize to a safe loader error
        raise ExtractionError(f"Failed to parse HTML ({type(exc).__name__})") from exc

    # Strip non-content elements so scripts/styles/metadata never reach the index.
    for tag in soup(["script", "style", "head"]):
        tag.decompose()

    text = soup.get_text(separator=" ", strip=True)
    return ExtractedDocument(text=text, page_count=1, parser="html")


def _load_docx(content: bytes) -> ExtractedDocument:
    """Extract paragraph text from a DOCX (OOXML) document via python-docx."""
    from docx import Document

    # Bound the inflated size BEFORE python-docx hands the XML to lxml.
    _guard_ooxml_zip_bomb(content, "DOCX")

    try:
        document = Document(io.BytesIO(content))
    except Exception as exc:  # noqa: BLE001 - normalize to a safe loader error
        raise ExtractionError(f"Failed to parse DOCX ({type(exc).__name__})") from exc

    paragraphs = [p.text for p in document.paragraphs if p.text and p.text.strip()]
    text = "\n".join(paragraphs)
    return ExtractedDocument(text=text, page_count=1, parser="docx")


def _load_pptx(content: bytes) -> ExtractedDocument:
    """Extract text from every slide of a PPTX deck via python-pptx.

    ``page_count`` is the slide count so the ``document_parsed`` event reports a
    meaningful page-equivalent for decks.
    """
    from pptx import Presentation

    # Bound the inflated size BEFORE python-pptx hands the XML to lxml.
    _guard_ooxml_zip_bomb(content, "PPTX")

    try:
        presentation = Presentation(io.BytesIO(content))
    except Exception as exc:  # noqa: BLE001 - normalize to a safe loader error
        raise ExtractionError(f"Failed to parse PPTX ({type(exc).__name__})") from exc

    slide_count = 0
    parts: list[str] = []
    for slide in presentation.slides:
        slide_count += 1
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text
                    if line and line.strip():
                        parts.append(line)
    text = "\n".join(parts)
    return ExtractedDocument(text=text, page_count=max(slide_count, 1), parser="pptx")


# ---------------------------------------------------------------------------
# Registry
# ---------------------------------------------------------------------------

# Extension → loader. The single source of truth for "what non-PDF formats does
# ingestion accept?" — referenced by the /upload route, the connectors, and
# upload() so the supported set is never duplicated.
_LOADERS: dict[str, Callable[[bytes], ExtractedDocument]] = {
    ".docx": _load_docx,
    ".pptx": _load_pptx,
    ".html": _load_html,
    ".htm": _load_html,
    ".md": _load_markdown,
    ".txt": _load_txt,
}

# PDF is handled by upload()'s existing PyMuPDF + Document Intelligence path, not
# by this registry. It is included in the supported set so callers can present
# one allowlist for "accepted by ingestion" routing decisions.
PDF_EXTENSION = ".pdf"

# Every extension ingestion accepts (PDF + registry formats), lowercased and
# dot-prefixed. The canonical allowlist for routing / filtering decisions.
SUPPORTED_EXTENSIONS: frozenset[str] = frozenset({PDF_EXTENSION, *_LOADERS})


def get_extension(filename: str) -> str:
    """Return the lowercased extension (with leading dot) of ``filename``.

    Returns an empty string when there is no extension. Uses the last dot so
    ``report.final.docx`` → ``.docx``.
    """
    name = (filename or "").rsplit("/", 1)[-1]
    dot = name.rfind(".")
    if dot <= 0:  # no dot, or a leading-dot dotfile with no real extension
        return ""
    return name[dot:].lower()


def is_supported_name(filename: str) -> bool:
    """True if ``filename``'s extension is an ingestion-supported format."""
    return get_extension(filename) in SUPPORTED_EXTENSIONS


def is_registry_format(filename: str) -> bool:
    """True if ``filename`` is a NON-PDF format handled by this registry.

    PDF returns False here because PDF text extraction stays in ``upload()``.
    """
    return get_extension(filename) in _LOADERS


def extract_text(filename: str, content: bytes) -> ExtractedDocument:
    """Dispatch by extension and extract plain text from ``content``.

    Raises:
        UnsupportedFormatError: the extension has no registered loader (callers
            should have pre-filtered; this is the defense-in-depth backstop).
        ExtractionError: a registered loader failed on malformed bytes.
    """
    ext = get_extension(filename)
    loader = _LOADERS.get(ext)
    if loader is None:
        # Do not echo the filename (may carry a path); the extension is safe.
        raise UnsupportedFormatError(f"Unsupported document format: {ext or '(no extension)'}")
    return loader(content)
